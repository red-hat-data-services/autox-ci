"""Document format handlers for dataset generation.

Supports saving documents in: txt, md, html, pptx, pdf, and images (png, jpg).
PDF and native image formats require binary data via save_binary_document().
"""

from __future__ import annotations

from pathlib import Path
from typing import Literal

DocumentFormat = Literal["txt", "md", "html", "pdf", "pptx", "png", "jpg"]


def save_document(
    content: str,
    output_path: Path,
    format: DocumentFormat = "txt",
    metadata: dict | None = None,
) -> Path:
    """Save document text content in the specified format.

    Args:
        content: Document text content
        output_path: Output path (extension will be replaced based on format)
        format: Output format - "txt" or "md"
        metadata: Optional metadata dict (used for markdown frontmatter)

    Returns:
        Path to the created file

    Note:
        This function is for text content only. For binary formats (PDF, PPTX, images),
        use save_binary_document() instead.
    """
    if format == "txt":
        return _save_txt(content, output_path)
    elif format == "md":
        return _save_markdown(content, output_path, metadata)
    elif format == "html":
        return _save_html(content, output_path, metadata)
    elif format == "pptx":
        return _save_pptx(content, output_path, metadata)
    elif format in ("pdf", "png", "jpg"):
        raise ValueError(
            f"Format '{format}' requires binary data. "
            f"Use save_binary_document() for {format} files, or use format='txt'/'md'/'html'/'pptx' for text content."
        )
    else:
        raise ValueError(f"Unsupported format: {format}. Use 'txt', 'md', 'html', or 'pptx'")


def _save_txt(content: str, output_path: Path) -> Path:
    """Save as plain text file."""
    # Don't use with_suffix() as it replaces everything after the last dot,
    # which truncates filenames like "doc_2410.14077v2" to "doc_2410.txt"
    # Instead, append .txt extension if not already present
    if output_path.suffix.lower() == ".txt":
        txt_path = output_path
    else:
        txt_path = Path(str(output_path) + ".txt")
    txt_path.write_text(content, encoding="utf-8")
    return txt_path


def _save_markdown(content: str, output_path: Path, metadata: dict | None = None) -> Path:
    """Save as markdown file with optional YAML frontmatter."""
    md_content = content

    if metadata:
        # Add YAML frontmatter
        frontmatter = "---\n"
        for key, value in metadata.items():
            # Escape special YAML characters in values
            value_str = str(value).replace('"', '\\"')
            frontmatter += f'{key}: "{value_str}"\n'
        frontmatter += "---\n\n"
        md_content = frontmatter + content

    # Don't use with_suffix() as it replaces everything after the last dot,
    # which truncates filenames like "doc_2410.14077v2" to "doc_2410.md"
    # Instead, append .md extension if not already present
    if output_path.suffix.lower() == ".md":
        md_path = output_path
    else:
        md_path = Path(str(output_path) + ".md")
    md_path.write_text(md_content, encoding="utf-8")
    return md_path


def _save_html(content: str, output_path: Path, metadata: dict | None = None) -> Path:
    """Save as HTML file."""
    title = (metadata or {}).get("title", "Document")
    paragraphs = "\n".join(f"<p>{p.strip()}</p>" for p in content.split("\n\n") if p.strip())
    html = (
        f"<!DOCTYPE html>\n<html><head><meta charset=\"utf-8\">"
        f"<title>{title}</title></head>\n<body>\n{paragraphs}\n</body></html>\n"
    )
    if output_path.suffix.lower() == ".html":
        html_path = output_path
    else:
        html_path = Path(str(output_path) + ".html")
    html_path.write_text(html, encoding="utf-8")
    return html_path


def _save_pptx(content: str, output_path: Path, metadata: dict | None = None) -> Path:
    """Save as PPTX with text distributed across slides."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_text = (metadata or {}).get("title", "Document")

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    txBox.text_frame.text = title_text
    txBox.text_frame.paragraphs[0].font.size = Pt(32)
    txBox.text_frame.paragraphs[0].font.bold = True

    # Content slides — split into ~500 char chunks by paragraph boundaries
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    chunk: list[str] = []
    chunk_len = 0
    for para in paragraphs:
        if chunk_len + len(para) > 500 and chunk:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = "\n\n".join(chunk)
            for p in tf.paragraphs:
                p.font.size = Pt(14)
            chunk = []
            chunk_len = 0
        chunk.append(para)
        chunk_len += len(para)
    if chunk:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = "\n\n".join(chunk)
        for p in tf.paragraphs:
            p.font.size = Pt(14)

    if output_path.suffix.lower() == ".pptx":
        pptx_path = output_path
    else:
        pptx_path = Path(str(output_path) + ".pptx")
    prs.save(str(pptx_path))
    return pptx_path


def save_binary_document(
    binary_data: bytes,
    output_path: Path,
    format: DocumentFormat,
) -> Path:
    """Save binary document data (PDF, PPTX, images).

    Args:
        binary_data: Binary file content
        output_path: Output path (extension will be appended based on format)
        format: Output format - "pdf", "pptx", "png", or "jpg"

    Returns:
        Path to the created file

    Raises:
        ValueError: If format is not a binary format
    """
    if format not in ("pdf", "pptx", "png", "jpg"):
        raise ValueError(
            f"Format '{format}' is not a binary format. "
            f"Use save_document() for text formats like 'txt' or 'md'."
        )

    # Don't use with_suffix() as it replaces everything after the last dot
    # Instead, append extension if not already present
    expected_suffix = f".{format}"
    if output_path.suffix.lower() == expected_suffix:
        output_file = output_path
    else:
        output_file = Path(str(output_path) + expected_suffix)
    output_file.write_bytes(binary_data)
    return output_file


def get_file_extension(format: DocumentFormat) -> str:
    """Get file extension for a given format.

    Args:
        format: Document format

    Returns:
        File extension including the dot (e.g., ".txt", ".md", ".pdf")
    """
    return f".{format}"


__all__ = [
    "DocumentFormat",
    "save_document",
    "save_binary_document",
    "get_file_extension",
]
